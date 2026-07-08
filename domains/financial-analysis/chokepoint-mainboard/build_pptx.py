#!/usr/bin/env python3
"""Canonical build_pptx.py · 3F Content F2 · 全屏图片嵌入 · 16:9 纯图无装饰

用法：
  1. 把 ImageGen 生成的 PNG 放到 images/ 目录，命名 p0_cover.png ... pN.png
  2. 修改 PAGES 列表
  3. 运行：python build_pptx.py

依赖：pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from pathlib import Path

# ===== 配置区 =====
IMAGES_DIR = Path(__file__).parent / "images"
OUTPUT_DIR = Path(__file__).parent / "output"
OUTPUT_FILE = OUTPUT_DIR / "presentation.pptx"

# 按页顺序列出图片文件名
PAGES = [
    "p0_cover.png",
    "p1_hook.png",
    "p2_mains.png",
    "p3_detail.png",
    "p4_detail2.png",
    "p5_counter.png",
    "p6_closing.png",
]

TITLE = "PPT 标题"
AUTHOR = "作者"
SUBJECT = "副标题/说明"
KEYWORDS = "关键词1 关键词2"
# ===================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.title = TITLE
prs.core_properties.author = AUTHOR
prs.core_properties.subject = SUBJECT
prs.core_properties.keywords = KEYWORDS

blank = prs.slide_layouts[6]
for fname in PAGES:
    s = prs.slides.add_slide(blank)
    path = IMAGES_DIR / fname
    if path.exists():
        s.shapes.add_picture(str(path), Emu(0), Emu(0),
                             width=prs.slide_width, height=prs.slide_height)
        print(f"  ✓ {fname}  ({path.stat().st_size // 1024} KB)")
    else:
        print(f"  ⚠ MISSING: {fname}")

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT_FILE))
mb = OUTPUT_FILE.stat().st_size / (1024 * 1024)
print(f"\n✅ {OUTPUT_FILE.name}  ·  {mb:.1f} MB  ·  {len(PAGES)} pages")
