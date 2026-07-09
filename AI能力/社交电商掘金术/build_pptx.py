#!/usr/bin/env python3
"""生成社交电商 × 两层拆解法 PPT 真实 .pptx"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu


PPT_DIR = Path("/Users/tingchi/.mavis/sessions/mvs_7e34de60b7c34f0e90937c2eaf2d24b9/workspace/社交电商掘金术-ppt")
IMAGES_DIR = PPT_DIR / "images"
OUTPUT_DIR = PPT_DIR / "output"
OUTPUT_PPTX = OUTPUT_DIR / "社交电商掘金术-cartman.pptx"

PAGES = [
    {"file": "p0_cover.png",       "title": "封面 · 社交电商 × 两层拆解法"},
    {"file": "p1_hook.png",        "title": "破局钩子 · 错误答案陷阱"},
    {"file": "p2_overview.png",    "title": "总纲 · 一张图理解"},
    {"file": "p3_layer1.png",      "title": "Part 1 · 第一层 四象限评估"},
    {"file": "p4_layer2.png",      "title": "Part 2 · 第二层 原子级 + 六维"},
    {"file": "p5_three_types.png", "title": "Step 6 · 3 种 Skill 类型"},
    {"file": "p6_roles.png",       "title": "4 角色升级"},
    {"file": "p7_takeaways.png",   "title": "Takeaways · 3 件 90 天"},
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    prs.core_properties.title = "社交电商 × 两层拆解法 · 8 页课件"
    prs.core_properties.author = "hmy1990116 · ai-training-methodology"
    prs.core_properties.subject = "社交电商 41 场景 / 两层拆解法 / 90 天落地 · 手绘叙事马卡龙"
    prs.core_properties.keywords = "社交电商 · 两层拆解法 · 四象限 · 原子级 · 六维 · Skill 类型 · 角色升级"

    blank_layout = prs.slide_layouts[6]

    for i, page in enumerate(PAGES, 1):
        slide = prs.slides.add_slide(blank_layout)
        img_path = IMAGES_DIR / page["file"]
        if not img_path.exists():
            print(f"  ⚠️ Missing: {img_path}")
            continue

        slide.shapes.add_picture(
            str(img_path),
            left=Emu(0), top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        size_kb = img_path.stat().st_size // 1024
        print(f"  ✓ Page {i}/8: {page['title']} ({size_kb} KB)")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))

    file_size_mb = OUTPUT_PPTX.stat().st_size / (1024 * 1024)
    print(f"\n{'='*60}")
    print(f"✅ .pptx 生成成功")
    print(f"   路径: {OUTPUT_PPTX}")
    print(f"   大小: {file_size_mb:.1f} MB")
    print(f"   页数: {len(PAGES)} · 16:9 宽屏")
    print(f"{'='*60}")


if __name__ == "__main__":
    build()
