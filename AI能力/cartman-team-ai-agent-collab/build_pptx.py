#!/usr/bin/env python3
"""生成真实 .pptx · 团队如何与 AI Agent 高效协作 · Cartman · 手绘叙事马卡龙

用户硬规则 (2026-07-03)：
- .pptx 不加任何画面装饰（页码 badge / 装饰元素都不要）
- 干净嵌入图片即可

输出：/Users/tingchi/.mavis/sessions/.../workspace/team-ai-agent-collab-ppt/output/team-ai-agent-collab-cartman.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu


PPT_DIR = Path("/Users/tingchi/.mavis/sessions/mvs_7e34de60b7c34f0e90937c2eaf2d24b9/workspace/team-ai-agent-collab-ppt")
IMAGES_DIR = PPT_DIR / "images"
OUTPUT_DIR = PPT_DIR / "output"
OUTPUT_PPTX = OUTPUT_DIR / "team-ai-agent-collab-cartman.pptx"

# 8 页 · 按 V2 范式 D「8 张黄金节奏」
PAGES = [
    {"file": "p0_cover.png",          "title": "封面 · 团队如何与 AI Agent 高效协作"},
    {"file": "p1_problem.png",        "title": "破局钩子 · 团队效能的 3 大症状"},
    {"file": "p2_vision.png",         "title": "总纲 · Context × Observation 双引擎"},
    {"file": "p3_context_layers.png", "title": "Part 1 · Context 三层办公室"},
    {"file": "p4_observation.png",    "title": "Part 2 · Observation 闭环"},
    {"file": "p5_memory_evolve.png",  "title": "Part 2 深化 · Memory Evolution"},
    {"file": "p6_human_role.png",     "title": "Part 3 · 4 角色升级"},
    {"file": "p7_takeaways.png",      "title": "Takeaways · 3 件今天开始"},
]


def build():
    prs = Presentation()

    # 16:9 宽屏
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 元数据修正（PowerPoint 识别关键）
    prs.core_properties.title = "团队如何与 AI Agent 高效协作"
    prs.core_properties.author = "Cartman · NODESK.AI STYLE"
    prs.core_properties.subject = "AI Native Team Collaboration · 手绘叙事马卡龙 · 8 页方法论精炼"
    prs.core_properties.keywords = "AI Agent · Context Engineering · Observation · Memory Evolution · Human Role"

    blank_layout = prs.slide_layouts[6]  # 完全空白版式

    for i, page in enumerate(PAGES, 1):
        slide = prs.slides.add_slide(blank_layout)
        img_path = IMAGES_DIR / page["file"]
        if not img_path.exists():
            print(f"  ⚠️ Missing: {img_path}")
            continue

        # 全屏图片，干净无装饰（按用户硬规则）
        slide.shapes.add_picture(
            str(img_path),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        size_kb = img_path.stat().st_size // 1024
        print(f"  ✓ Page {i}/8: {page['title']} ({size_kb} KB)")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))

    file_size_mb = OUTPUT_PPTX.stat().st_size / (1024 * 1024)
    print(f"\n{'='*60}")
    print(f"✅ .pptx 生成成功")
    print(f"   路径: {OUTPUT_PPTX}")
    print(f"   大小: {file_size_mb:.1f} MB")
    print(f"   页数: {len(PAGES)} · 16:9 宽屏")
    print(f"   风格: 手绘叙事 + 马卡龙经典")
    print(f"   装饰: 无（用户偏好）")
    print(f"{'='*60}")


if __name__ == "__main__":
    build()
