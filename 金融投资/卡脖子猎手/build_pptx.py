#!/usr/bin/env python3
"""卡脖子选股 PPTX · 3F Content F2 · 全屏图片嵌入 · 16:9 纯图无装饰
用法：python build_pptx.py
前提：images/ 下有 p0_cover.png ~ p8_action.png
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from pathlib import Path

IMAGES_DIR = Path(__file__).parent / "images"
OUTPUT_DIR = Path(__file__).parent / "output"
OUTPUT_FILE = OUTPUT_DIR / "chokepoint_mainboard.pptx"

PAGES = [
    ("p0_cover.png",   "P0 · 封面 · 卡脖子选股 · 主板专版"),
    ("p1_hook.png",    "P1 · 破局钩子 · 反制牌"),
    ("p2_mains.png",   "P2 · 三大主线 · 光刻胶/稀土/前驱体"),
    ("p3_shenghe.png", "P3 · S级 · 盛和资源 600392"),
    ("p4_tongcheng.png","P4 · A级 · 彤程新材 603650"),
    ("p5_yake.png",    "P5 · A级 · 雅克科技 002409"),
    ("p6_youyan.png",  "P6 · B级 · 有研新材 600206（排除）"),
    ("p7_compare.png", "P7 · 四维信号对比表"),
    ("p8_action.png",  "P8 · 操作建议 + 仓位纪律"),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.title = "卡脖子选股报告 · 主板专版 · 2026-07-07"
prs.core_properties.author = "WorkBuddy AI · stock_wisdom · ImageGen"
prs.core_properties.subject = "Serenity 卡脖子六步法选股 · 深色科技风 · 3F Content"
prs.core_properties.keywords = "卡脖子 选股 主板 盛和资源 彤程新材 雅克科技 稀土 光刻胶"

blank = prs.slide_layouts[6]
ok = 0
for fname, title in PAGES:
    s = prs.slides.add_slide(blank)
    path = IMAGES_DIR / fname
    if path.exists():
        s.shapes.add_picture(str(path), Emu(0), Emu(0),
                             width=prs.slide_width, height=prs.slide_height)
        kb = path.stat().st_size // 1024
        print(f"  ✓ {title} ({kb} KB)")
        ok += 1
    else:
        print(f"  ⚠ MISSING: {fname} — {title}")

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT_FILE))
mb = OUTPUT_FILE.stat().st_size / (1024 * 1024)
print(f"\n✅ {OUTPUT_FILE.name}  ·  {mb:.1f} MB  ·  {ok}/{len(PAGES)} pages")
if ok != len(PAGES):
    print("⚠ Not all images present — run ImageGen for missing pages first")
